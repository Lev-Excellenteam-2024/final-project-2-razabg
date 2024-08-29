from pptx import Presentation


class PresentationReader:
    """
    A class to read and extract text from a PowerPoint presentation (.pptx).

    Attributes:
        file_path (str): The path to the PowerPoint presentation file.
        presentation (Presentation): An instance of the Presentation class from the pptx library.
        presentation_content (dict): A dictionary storing the text content of each slide,
                                      where keys are slide numbers and values are concatenated text strings.

    Methods:
        extract_text():
            Extracts all text from the presentation and stores it in a dictionary.
            Each entry in the dictionary corresponds to a slide number with its text content.
            Returns:
                dict: A dictionary where the keys are slide numbers (int) and the values are strings
                      containing the text from each slide.
    """

    def __init__(self, file_path):
        """
        Initializes the PresentationReader with the path to a PowerPoint presentation.

        Args:
            file_path (str): The path to the PowerPoint presentation file.
        """
        self.file_path = file_path
        self.presentation = Presentation(file_path)
        self.presentation_content = {}

    def extract_text(self):
        """
        Extracts all text from the presentation and stores it in the presentation_content dictionary.

        Iterates through each slide and each shape within the slide, collecting text from shapes
        that contain text. The text is concatenated and stored in the presentation_content dictionary
        with the slide number as the key. Empty slides are skipped.

        Returns:
            dict: A dictionary where the keys are slide numbers (int) and the values are strings
                  containing the text from each slide.
        """
        separator: str = " "
        for slide_num, slide in enumerate(self.presentation.slides, start=1):
            slide_text = []

            # Iterate through each shape in the slide
            for shape in slide.shapes:
                # Check if the shape has text (e.g., text boxes, titles, etc.)
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            # Join the collected text for this slide and add it to presentation_content if not empty
            if slide_text:
                self.presentation_content[slide_num] = separator.join(slide_text)

        return self.presentation_content
